import os
import json
import logging
import re
from typing import Dict, Any, List
from django.conf import settings
from pptx import Presentation
from pypdf import PdfReader
from apps.common.llm import get_text_llm

logger = logging.getLogger(__name__)

class DocProcessor:
    def extract_text_from_pdf(self, file_path: str) -> str:
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""

    def extract_text_from_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return ""

    def analyze_content(self, file_path: str, file_type: str) -> Dict[str, Any]:
        # 1. 提取文本
        if 'pdf' in file_type.lower():
            content = self.extract_text_from_pdf(file_path)
        elif 'ppt' in file_type.lower():
            content = self.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

        print(f"[DocProcessor] 成功从文档提取内容，长度: {len(content)} 字")

        if not content.strip():
            raise ValueError("未能从文档中提取到有效文本内容。")

        # 2. 调用 LLM 进行分析
        llm = get_text_llm()
        
        prompt = f"""
你是一个专业的消防安全培训专家。请分析以下从【{file_type}】文档中提取的内容，并生成一套完整的培训教学大纲和配套的实战演练案例。

文档内容：
{content[:8000]}

请严格按以下 JSON 格式输出，不要包含任何 markdown 代码块标记或额外解释：
{{
  "teaching_content": "一份系统化的教学讲义，包含知识点讲解，适合数字人播报，约 500-800 字。使用 Markdown 格式。",
  "scenarios": [
    {{
      "title": "场景标题",
      "topic": "主题分类",
      "difficulty": "low",
      "description": "事故现场的详细描述",
      "correct_actions": "标准的应急处置流程步骤",
      "analysis": "演练后的点评与原因分析总结"
    }}
  ]
}}
要求生成至少 3 个不同难度的演练场景（难度可选 low, medium, high）。
"""
        
        raw_resp = llm.chat([{"role": "user", "content": prompt}], temperature=0.3, max_tokens=3500)
        
        # 尝试解析 JSON
        try:
            # 简单清洗
            clean_json = raw_resp.strip()
            # 去除可能存在的 Markdown 代码块标记
            clean_json = re.sub(r'^```[a-z]*\s*|\s*```$', '', clean_json, flags=re.MULTILINE | re.IGNORECASE)
            
            # 找到第一个 { 和最后一个 }
            start = clean_json.find('{')
            end = clean_json.rfind('}')
            if start != -1 and end != -1:
                clean_json = clean_json[start:end+1]
                
            # 使用 strict=False 允许控制字符(如换行符)存在于字符串中
            result = json.loads(clean_json, strict=False)
            
            if 'teaching_content' not in result or 'scenarios' not in result:
                # 尝试检查 key 是否被加上了额外引号或空格
                # 兼容性处理
                keys = result.keys()
                t_key = next((k for k in keys if 'teaching' in k), None)
                s_key = next((k for k in keys if 'scenario' in k), None)
                if t_key and s_key:
                    result = {
                        'teaching_content': result[t_key],
                        'scenarios': result[s_key]
                    }
                else:
                    raise ValueError("输出结果缺少必要字段: teaching_content, scenarios")
                    
            return result
        except Exception as e:
            logger.error(f"Failed to parse LLM response: {e}\nResponse: {raw_resp}")
            # 最后的挣扎：如果解析失败，尝试直接捕获可能存在的文本片段（针对某些极其不规范的输出）
            raise ValueError(f"大模型返回结果解析失败: {str(e)}")
